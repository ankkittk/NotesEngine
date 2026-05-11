import os

import fitz
from docx import Document
from pptx import Presentation

from ..core.config import (
    OCR_MIN_TEXT_THRESHOLD,
    VISION_BATCH_SIZE,
    VISION_PDF_RENDER_SCALE,
)
from .vision.vision_extractor import extract_batch_images


def pdf_page_has_image(page):
    return len(page.get_images(full=True)) > 0


def should_use_vision(text, has_image):
    return has_image and len(text.split()) < OCR_MIN_TEXT_THRESHOLD


def _make_record(text, source, page):
    return {
        "text": text or "",
        "source": source,
        "page": page,
    }


def load_pdf(file_path):
    source = os.path.basename(file_path)
    doc = fitz.open(file_path)

    results = []
    batch_imgs = []
    batch_records = []

    try:
        for i, page in enumerate(doc, start=1):
            text = page.get_text()
            has_image = pdf_page_has_image(page)
            record = _make_record(text, source, i)

            if should_use_vision(text, has_image):
                pix = page.get_pixmap(
                    matrix=fitz.Matrix(VISION_PDF_RENDER_SCALE, VISION_PDF_RENDER_SCALE)
                )
                batch_imgs.append(pix.tobytes("png"))
                batch_records.append(record)
                results.append(record)

                if len(batch_imgs) == VISION_BATCH_SIZE:
                    outs = extract_batch_images(batch_imgs)
                    for rec, out in zip(batch_records, outs):
                        rec["text"] = out or ""
                    batch_imgs, batch_records = [], []
            else:
                results.append(record)

        if batch_imgs:
            print(f"   → Vision batch ({len(batch_imgs)} pages)...", flush=True)
            outs = extract_batch_images(batch_imgs)
            for rec, out in zip(batch_records, outs):
                rec["text"] = out or ""

    finally:
        doc.close()

    return [r for r in results if r.get("text", "").strip()]


def load_txt(file_path):
    source = os.path.basename(file_path)
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()

    return [
        {
            "text": text,
            "source": source,
            "page": 1,
        }
    ] if text.strip() else []


def load_docx(file_path):
    source = os.path.basename(file_path)
    doc = Document(file_path)
    text = "\n".join(p.text for p in doc.paragraphs).strip()

    return [
        {
            "text": text,
            "source": source,
            "page": 1,
        }
    ] if text else []


def load_pptx(file_path):
    source = os.path.basename(file_path)
    prs = Presentation(file_path)
    records = []

    for idx, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)

        text = "\n".join(parts).strip()
        if text:
            records.append(
                {
                    "text": text,
                    "source": source,
                    "page": idx,
                }
            )

    return records


def load_document(file_path):
    ext = os.path.splitext(file_path.lower())[1]

    if ext == ".pdf":
        return load_pdf(file_path)
    if ext == ".txt":
        return load_txt(file_path)
    if ext == ".docx":
        return load_docx(file_path)
    if ext == ".pptx":
        return load_pptx(file_path)

    return []
