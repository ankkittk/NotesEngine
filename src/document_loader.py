import os
import fitz
from docx import Document
from pptx import Presentation

from vision_extractor import extract_batch_images

BATCH_SIZE = 4


def pdf_page_has_image(page):
    return len(page.get_images(full=True)) > 0


def should_use_vision(text, has_image):
    # minimal gate to avoid over-calling API
    return has_image and len(text.split()) < 30


def load_pdf(file_path):
    doc = fitz.open(file_path)

    results = []
    batch_imgs = []
    batch_indices = []

    for i, page in enumerate(doc):
        text = page.get_text()
        has_image = pdf_page_has_image(page)

        if should_use_vision(text, has_image):
            pix = page.get_pixmap()
            batch_imgs.append(pix.tobytes("png"))
            batch_indices.append(i)
            results.append(None)  # placeholder

            if len(batch_imgs) == BATCH_SIZE:
                outs = extract_batch_images(batch_imgs)
                for idx, out in zip(batch_indices, outs):
                    results[idx] = out
                batch_imgs, batch_indices = [], []
        else:
            results.append(text)

    # flush remaining batch
    if batch_imgs:
        outs = extract_batch_images(batch_imgs)
        for idx, out in zip(batch_indices, outs):
            results[idx] = out

    return "\n".join(r for r in results if r)


def load_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_docx(file_path):
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def load_pptx(file_path):
    prs = Presentation(file_path)
    slides = []

    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
        slides.append("\n".join(parts))

    return "\n".join(slides)


def load_document(file_path):
    ext = os.path.splitext(file_path.lower())[1]

    if ext == ".pdf":
        return load_pdf(file_path)
    if ext == ".txt":
        return load_txt(file_path)
    if ext == ".docx":
        return load_docx(file_path)
    if ext == ".pptx":
        return load_pptx(file_path)

    return ""
